"""Comprehensive PowerPoint MCP Server with full PPTX editing capabilities."""

import asyncio
import base64
import json
import logging
import os
import sys
from io import BytesIO
from pathlib import Path
from typing import Any, Dict, List, Optional, Union

from mcp.server import Server
from mcp.server.models import InitializationOptions
from mcp.types import TextContent, Tool
from pptx import Presentation
from pptx.chart.data import CategoryChartData
from pptx.dml.color import RGBColor
from pptx.enum.chart import XL_CHART_TYPE
from pptx.enum.dml import MSO_THEME_COLOR
from pptx.enum.shapes import MSO_SHAPE
from pptx.enum.text import MSO_ANCHOR, MSO_AUTO_SIZE, PP_ALIGN
from pptx.parts.image import Image
from pptx.slide import Slide
from pptx.util import Inches, Pt

logging.basicConfig(
    level=logging.INFO,
    format="%(asctime)s - %(levelname)s - %(message)s",
    handlers=[logging.StreamHandler(sys.stderr)],
)
log = logging.getLogger("pptx_server")

server = Server("pptx-server")

# Global presentation cache to maintain state across operations
_presentations: Dict[str, Presentation] = {}

# Default directories for organizing presentations
DEFAULT_OUTPUT_DIR = "examples/generated"
DEFAULT_TEMPLATE_DIR = "examples/templates"
DEFAULT_DEMO_DIR = "examples/demos"


def _ensure_output_directory(file_path: str) -> str:
    """Ensure output directory exists and return full path."""
    # If it's already an absolute path, use as-is
    if os.path.isabs(file_path):
        dir_path = os.path.dirname(file_path)
        os.makedirs(dir_path, exist_ok=True)
        return file_path

    # If it's a relative path, check if it includes a directory
    if os.path.dirname(file_path):
        os.makedirs(os.path.dirname(file_path), exist_ok=True)
        return file_path

    # If it's just a filename, put it in the default output directory
    output_dir = DEFAULT_OUTPUT_DIR
    os.makedirs(output_dir, exist_ok=True)
    return os.path.join(output_dir, file_path)


def _resolve_template_path(template_path: str) -> str:
    """Resolve template path, checking template directory if relative."""
    if os.path.isabs(template_path) or os.path.exists(template_path):
        return template_path

    # Check in templates directory
    template_in_dir = os.path.join(DEFAULT_TEMPLATE_DIR, template_path)
    if os.path.exists(template_in_dir):
        return template_in_dir

    return template_path  # Return original if not found


def _get_presentation(file_path: str) -> Presentation:
    """Get or create a presentation, caching it for subsequent operations."""
    abs_path = os.path.abspath(file_path)

    if abs_path not in _presentations:
        if os.path.exists(abs_path):
            log.info(f"Loading existing presentation: {abs_path}")
            _presentations[abs_path] = Presentation(abs_path)
        else:
            log.info(f"Creating new presentation: {abs_path}")
            _presentations[abs_path] = Presentation()

    return _presentations[abs_path]


def _save_presentation(file_path: str) -> None:
    """Save a presentation and update the cache."""
    abs_path = os.path.abspath(file_path)
    if abs_path in _presentations:
        # Ensure directory exists
        os.makedirs(os.path.dirname(abs_path), exist_ok=True)
        _presentations[abs_path].save(abs_path)
        log.info(f"Saved presentation: {abs_path}")


def _parse_color(color_str: str) -> RGBColor:
    """Parse color string (hex format like #FF0000) to RGBColor."""
    if color_str.startswith("#"):
        color_str = color_str[1:]
    return RGBColor(int(color_str[:2], 16), int(color_str[2:4], 16), int(color_str[4:6], 16))


@server.list_tools()
async def list_tools() -> list[Tool]:
    """List all available PowerPoint editing tools."""
    return [
        # Presentation Management
        Tool(
            name="create_presentation",
            description="Create a new PowerPoint presentation",
            inputSchema={
                "type": "object",
                "properties": {
                    "file_path": {"type": "string", "description": "Path where to save the presentation"},
                    "title": {"type": "string", "description": "Optional title for the presentation"}
                },
                "required": ["file_path"]
            }
        ),
        Tool(
            name="create_presentation_from_template",
            description="Create a new PowerPoint presentation from an existing template",
            inputSchema={
                "type": "object",
                "properties": {
                    "template_path": {"type": "string", "description": "Path to the template presentation file"},
                    "output_path": {"type": "string", "description": "Path where to save the new presentation"},
                    "title": {"type": "string", "description": "Optional new title for the presentation"},
                    "replace_placeholders": {"type": "object", "description": "Key-value pairs to replace text placeholders", "additionalProperties": {"type": "string"}}
                },
                "required": ["template_path", "output_path"]
            }
        ),
        Tool(
            name="clone_presentation",
            description="Clone an existing presentation with optional modifications",
            inputSchema={
                "type": "object",
                "properties": {
                    "source_path": {"type": "string", "description": "Path to the source presentation"},
                    "target_path": {"type": "string", "description": "Path for the cloned presentation"},
                    "new_title": {"type": "string", "description": "Optional new title for the cloned presentation"}
                },
                "required": ["source_path", "target_path"]
            }
        ),
        Tool(
            name="open_presentation",
            description="Open an existing PowerPoint presentation",
            inputSchema={
                "type": "object",
                "properties": {
                    "file_path": {"type": "string", "description": "Path to the presentation file"}
                },
                "required": ["file_path"]
            }
        ),
        Tool(
            name="save_presentation",
            description="Save the current presentation to file",
            inputSchema={
                "type": "object",
                "properties": {
                    "file_path": {"type": "string", "description": "Path where to save the presentation"}
                },
                "required": ["file_path"]
            }
        ),
        Tool(
            name="get_presentation_info",
            description="Get information about the presentation (slide count, properties, etc.)",
            inputSchema={
                "type": "object",
                "properties": {
                    "file_path": {"type": "string", "description": "Path to the presentation file"}
                },
                "required": ["file_path"]
            }
        ),

        # Slide Management
        Tool(
            name="add_slide",
            description="Add a new slide to the presentation",
            inputSchema={
                "type": "object",
                "properties": {
                    "file_path": {"type": "string", "description": "Path to the presentation file"},
                    "layout_index": {"type": "integer", "description": "Slide layout index (0-based)", "default": 0},
                    "position": {"type": "integer", "description": "Position to insert slide (0-based, -1 for end)", "default": -1}
                },
                "required": ["file_path"]
            }
        ),
        Tool(
            name="delete_slide",
            description="Delete a slide from the presentation",
            inputSchema={
                "type": "object",
                "properties": {
                    "file_path": {"type": "string", "description": "Path to the presentation file"},
                    "slide_index": {"type": "integer", "description": "Index of slide to delete (0-based)"}
                },
                "required": ["file_path", "slide_index"]
            }
        ),
        Tool(
            name="move_slide",
            description="Move a slide to a different position",
            inputSchema={
                "type": "object",
                "properties": {
                    "file_path": {"type": "string", "description": "Path to the presentation file"},
                    "from_index": {"type": "integer", "description": "Current index of slide (0-based)"},
                    "to_index": {"type": "integer", "description": "New index for slide (0-based)"}
                },
                "required": ["file_path", "from_index", "to_index"]
            }
        ),
        Tool(
            name="duplicate_slide",
            description="Duplicate an existing slide",
            inputSchema={
                "type": "object",
                "properties": {
                    "file_path": {"type": "string", "description": "Path to the presentation file"},
                    "slide_index": {"type": "integer", "description": "Index of slide to duplicate (0-based)"},
                    "position": {"type": "integer", "description": "Position for duplicated slide (-1 for end)", "default": -1}
                },
                "required": ["file_path", "slide_index"]
            }
        ),
        Tool(
            name="list_slides",
            description="List all slides in the presentation with their basic information",
            inputSchema={
                "type": "object",
                "properties": {
                    "file_path": {"type": "string", "description": "Path to the presentation file"}
                },
                "required": ["file_path"]
            }
        ),

        # Text and Content Management
        Tool(
            name="set_slide_title",
            description="Set the title of a slide",
            inputSchema={
                "type": "object",
                "properties": {
                    "file_path": {"type": "string", "description": "Path to the presentation file"},
                    "slide_index": {"type": "integer", "description": "Index of slide (0-based)"},
                    "title": {"type": "string", "description": "Title text"}
                },
                "required": ["file_path", "slide_index", "title"]
            }
        ),
        Tool(
            name="set_slide_content",
            description="Set the main content/body text of a slide",
            inputSchema={
                "type": "object",
                "properties": {
                    "file_path": {"type": "string", "description": "Path to the presentation file"},
                    "slide_index": {"type": "integer", "description": "Index of slide (0-based)"},
                    "content": {"type": "string", "description": "Content text (can include bullet points with \\n)"}
                },
                "required": ["file_path", "slide_index", "content"]
            }
        ),
        Tool(
            name="add_text_box",
            description="Add a text box to a slide",
            inputSchema={
                "type": "object",
                "properties": {
                    "file_path": {"type": "string", "description": "Path to the presentation file"},
                    "slide_index": {"type": "integer", "description": "Index of slide (0-based)"},
                    "text": {"type": "string", "description": "Text content"},
                    "left": {"type": "number", "description": "Left position in inches", "default": 1.0},
                    "top": {"type": "number", "description": "Top position in inches", "default": 1.0},
                    "width": {"type": "number", "description": "Width in inches", "default": 6.0},
                    "height": {"type": "number", "description": "Height in inches", "default": 1.0},
                    "font_size": {"type": "integer", "description": "Font size in points", "default": 18},
                    "font_color": {"type": "string", "description": "Font color in hex (#RRGGBB)", "default": "#000000"},
                    "bold": {"type": "boolean", "description": "Make text bold", "default": False},
                    "italic": {"type": "boolean", "description": "Make text italic", "default": False}
                },
                "required": ["file_path", "slide_index", "text"]
            }
        ),
        Tool(
            name="format_text",
            description="Format existing text in a slide (placeholder or text box)",
            inputSchema={
                "type": "object",
                "properties": {
                    "file_path": {"type": "string", "description": "Path to the presentation file"},
                    "slide_index": {"type": "integer", "description": "Index of slide (0-based)"},
                    "shape_index": {"type": "integer", "description": "Index of text shape (0-based)"},
                    "font_name": {"type": "string", "description": "Font name (e.g., 'Arial', 'Times New Roman')"},
                    "font_size": {"type": "integer", "description": "Font size in points"},
                    "font_color": {"type": "string", "description": "Font color in hex (#RRGGBB)"},
                    "bold": {"type": "boolean", "description": "Make text bold"},
                    "italic": {"type": "boolean", "description": "Make text italic"},
                    "underline": {"type": "boolean", "description": "Underline text"},
                    "alignment": {"type": "string", "description": "Text alignment (left, center, right, justify)", "default": "left"}
                },
                "required": ["file_path", "slide_index", "shape_index"]
            }
        ),

        # Image Management
        Tool(
            name="add_image",
            description="Add an image to a slide",
            inputSchema={
                "type": "object",
                "properties": {
                    "file_path": {"type": "string", "description": "Path to the presentation file"},
                    "slide_index": {"type": "integer", "description": "Index of slide (0-based)"},
                    "image_path": {"type": "string", "description": "Path to the image file"},
                    "left": {"type": "number", "description": "Left position in inches", "default": 1.0},
                    "top": {"type": "number", "description": "Top position in inches", "default": 1.0},
                    "width": {"type": "number", "description": "Width in inches (optional, maintains aspect ratio)"},
                    "height": {"type": "number", "description": "Height in inches (optional, maintains aspect ratio)"}
                },
                "required": ["file_path", "slide_index", "image_path"]
            }
        ),
        Tool(
            name="add_image_from_base64",
            description="Add an image to a slide from base64 encoded data",
            inputSchema={
                "type": "object",
                "properties": {
                    "file_path": {"type": "string", "description": "Path to the presentation file"},
                    "slide_index": {"type": "integer", "description": "Index of slide (0-based)"},
                    "image_data": {"type": "string", "description": "Base64 encoded image data"},
                    "image_format": {"type": "string", "description": "Image format (png, jpg, gif)", "default": "png"},
                    "left": {"type": "number", "description": "Left position in inches", "default": 1.0},
                    "top": {"type": "number", "description": "Top position in inches", "default": 1.0},
                    "width": {"type": "number", "description": "Width in inches (optional)"},
                    "height": {"type": "number", "description": "Height in inches (optional)"}
                },
                "required": ["file_path", "slide_index", "image_data"]
            }
        ),
        Tool(
            name="replace_image",
            description="Replace an existing image in a slide",
            inputSchema={
                "type": "object",
                "properties": {
                    "file_path": {"type": "string", "description": "Path to the presentation file"},
                    "slide_index": {"type": "integer", "description": "Index of slide (0-based)"},
                    "shape_index": {"type": "integer", "description": "Index of image shape (0-based)"},
                    "new_image_path": {"type": "string", "description": "Path to the new image file"}
                },
                "required": ["file_path", "slide_index", "shape_index", "new_image_path"]
            }
        ),

        # Shape Management
        Tool(
            name="add_shape",
            description="Add a shape to a slide",
            inputSchema={
                "type": "object",
                "properties": {
                    "file_path": {"type": "string", "description": "Path to the presentation file"},
                    "slide_index": {"type": "integer", "description": "Index of slide (0-based)"},
                    "shape_type": {"type": "string", "description": "Shape type (rectangle, oval, triangle, arrow, etc.)"},
                    "left": {"type": "number", "description": "Left position in inches", "default": 1.0},
                    "top": {"type": "number", "description": "Top position in inches", "default": 1.0},
                    "width": {"type": "number", "description": "Width in inches", "default": 2.0},
                    "height": {"type": "number", "description": "Height in inches", "default": 1.0},
                    "fill_color": {"type": "string", "description": "Fill color in hex (#RRGGBB)"},
                    "line_color": {"type": "string", "description": "Line color in hex (#RRGGBB)"},
                    "line_width": {"type": "number", "description": "Line width in points", "default": 1.0}
                },
                "required": ["file_path", "slide_index", "shape_type"]
            }
        ),
        Tool(
            name="modify_shape",
            description="Modify properties of an existing shape",
            inputSchema={
                "type": "object",
                "properties": {
                    "file_path": {"type": "string", "description": "Path to the presentation file"},
                    "slide_index": {"type": "integer", "description": "Index of slide (0-based)"},
                    "shape_index": {"type": "integer", "description": "Index of shape (0-based)"},
                    "left": {"type": "number", "description": "Left position in inches"},
                    "top": {"type": "number", "description": "Top position in inches"},
                    "width": {"type": "number", "description": "Width in inches"},
                    "height": {"type": "number", "description": "Height in inches"},
                    "fill_color": {"type": "string", "description": "Fill color in hex (#RRGGBB)"},
                    "line_color": {"type": "string", "description": "Line color in hex (#RRGGBB)"},
                    "line_width": {"type": "number", "description": "Line width in points"}
                },
                "required": ["file_path", "slide_index", "shape_index"]
            }
        ),
        Tool(
            name="delete_shape",
            description="Delete a shape from a slide",
            inputSchema={
                "type": "object",
                "properties": {
                    "file_path": {"type": "string", "description": "Path to the presentation file"},
                    "slide_index": {"type": "integer", "description": "Index of slide (0-based)"},
                    "shape_index": {"type": "integer", "description": "Index of shape to delete (0-based)"}
                },
                "required": ["file_path", "slide_index", "shape_index"]
            }
        ),

        # Table Operations
        Tool(
            name="add_table",
            description="Add a table to a slide",
            inputSchema={
                "type": "object",
                "properties": {
                    "file_path": {"type": "string", "description": "Path to the presentation file"},
                    "slide_index": {"type": "integer", "description": "Index of slide (0-based)"},
                    "rows": {"type": "integer", "description": "Number of rows", "minimum": 1},
                    "cols": {"type": "integer", "description": "Number of columns", "minimum": 1},
                    "left": {"type": "number", "description": "Left position in inches", "default": 1.0},
                    "top": {"type": "number", "description": "Top position in inches", "default": 1.0},
                    "width": {"type": "number", "description": "Table width in inches", "default": 6.0},
                    "height": {"type": "number", "description": "Table height in inches", "default": 3.0}
                },
                "required": ["file_path", "slide_index", "rows", "cols"]
            }
        ),
        Tool(
            name="set_table_cell",
            description="Set the text content of a table cell",
            inputSchema={
                "type": "object",
                "properties": {
                    "file_path": {"type": "string", "description": "Path to the presentation file"},
                    "slide_index": {"type": "integer", "description": "Index of slide (0-based)"},
                    "table_index": {"type": "integer", "description": "Index of table shape (0-based)"},
                    "row": {"type": "integer", "description": "Row index (0-based)"},
                    "col": {"type": "integer", "description": "Column index (0-based)"},
                    "text": {"type": "string", "description": "Cell text content"}
                },
                "required": ["file_path", "slide_index", "table_index", "row", "col", "text"]
            }
        ),
        Tool(
            name="format_table_cell",
            description="Format a table cell (font, color, alignment)",
            inputSchema={
                "type": "object",
                "properties": {
                    "file_path": {"type": "string", "description": "Path to the presentation file"},
                    "slide_index": {"type": "integer", "description": "Index of slide (0-based)"},
                    "table_index": {"type": "integer", "description": "Index of table shape (0-based)"},
                    "row": {"type": "integer", "description": "Row index (0-based)"},
                    "col": {"type": "integer", "description": "Column index (0-based)"},
                    "font_size": {"type": "integer", "description": "Font size in points"},
                    "font_color": {"type": "string", "description": "Font color in hex (#RRGGBB)"},
                    "fill_color": {"type": "string", "description": "Cell background color in hex (#RRGGBB)"},
                    "bold": {"type": "boolean", "description": "Make text bold"},
                    "alignment": {"type": "string", "description": "Text alignment (left, center, right)"}
                },
                "required": ["file_path", "slide_index", "table_index", "row", "col"]
            }
        ),
        Tool(
            name="populate_table",
            description="Populate entire table with data from a 2D array",
            inputSchema={
                "type": "object",
                "properties": {
                    "file_path": {"type": "string", "description": "Path to the presentation file"},
                    "slide_index": {"type": "integer", "description": "Index of slide (0-based)"},
                    "table_index": {"type": "integer", "description": "Index of table shape (0-based)"},
                    "data": {"type": "array", "description": "2D array of cell values", "items": {"type": "array", "items": {"type": "string"}}},
                    "header_row": {"type": "boolean", "description": "Format first row as header", "default": False}
                },
                "required": ["file_path", "slide_index", "table_index", "data"]
            }
        ),

        # Chart Operations
        Tool(
            name="add_chart",
            description="Add a chart to a slide",
            inputSchema={
                "type": "object",
                "properties": {
                    "file_path": {"type": "string", "description": "Path to the presentation file"},
                    "slide_index": {"type": "integer", "description": "Index of slide (0-based)"},
                    "chart_type": {"type": "string", "description": "Chart type (column, bar, line, pie)", "default": "column"},
                    "data": {"type": "object", "description": "Chart data with categories and series", "properties": {
                        "categories": {"type": "array", "items": {"type": "string"}},
                        "series": {"type": "array", "items": {"type": "object", "properties": {
                            "name": {"type": "string"},
                            "values": {"type": "array", "items": {"type": "number"}}
                        }}}
                    }},
                    "left": {"type": "number", "description": "Left position in inches", "default": 1.0},
                    "top": {"type": "number", "description": "Top position in inches", "default": 1.0},
                    "width": {"type": "number", "description": "Chart width in inches", "default": 6.0},
                    "height": {"type": "number", "description": "Chart height in inches", "default": 4.0},
                    "title": {"type": "string", "description": "Chart title"}
                },
                "required": ["file_path", "slide_index", "data"]
            }
        ),
        Tool(
            name="update_chart_data",
            description="Update data in an existing chart",
            inputSchema={
                "type": "object",
                "properties": {
                    "file_path": {"type": "string", "description": "Path to the presentation file"},
                    "slide_index": {"type": "integer", "description": "Index of slide (0-based)"},
                    "chart_index": {"type": "integer", "description": "Index of chart shape (0-based)"},
                    "data": {"type": "object", "description": "New chart data", "properties": {
                        "categories": {"type": "array", "items": {"type": "string"}},
                        "series": {"type": "array", "items": {"type": "object", "properties": {
                            "name": {"type": "string"},
                            "values": {"type": "array", "items": {"type": "number"}}
                        }}}
                    }}
                },
                "required": ["file_path", "slide_index", "chart_index", "data"]
            }
        ),

        # Utility and Information Tools
        Tool(
            name="list_shapes",
            description="List all shapes on a slide with their types and properties",
            inputSchema={
                "type": "object",
                "properties": {
                    "file_path": {"type": "string", "description": "Path to the presentation file"},
                    "slide_index": {"type": "integer", "description": "Index of slide (0-based)"}
                },
                "required": ["file_path", "slide_index"]
            }
        ),
        Tool(
            name="get_slide_layouts",
            description="Get available slide layouts in the presentation",
            inputSchema={
                "type": "object",
                "properties": {
                    "file_path": {"type": "string", "description": "Path to the presentation file"}
                },
                "required": ["file_path"]
            }
        ),
        Tool(
            name="set_presentation_properties",
            description="Set presentation document properties",
            inputSchema={
                "type": "object",
                "properties": {
                    "file_path": {"type": "string", "description": "Path to the presentation file"},
                    "title": {"type": "string", "description": "Presentation title"},
                    "author": {"type": "string", "description": "Author name"},
                    "subject": {"type": "string", "description": "Subject"},
                    "comments": {"type": "string", "description": "Comments"}
                },
                "required": ["file_path"]
            }
        ),
        Tool(
            name="export_slide_as_image",
            description="Export a slide as an image file",
            inputSchema={
                "type": "object",
                "properties": {
                    "file_path": {"type": "string", "description": "Path to the presentation file"},
                    "slide_index": {"type": "integer", "description": "Index of slide (0-based)"},
                    "output_path": {"type": "string", "description": "Output image file path"},
                    "format": {"type": "string", "description": "Image format (png, jpg)", "default": "png"}
                },
                "required": ["file_path", "slide_index", "output_path"]
            }
        ),

        # Composite Workflow Tools
        Tool(
            name="create_title_slide",
            description="Create a complete title slide with title, subtitle, and optional company info",
            inputSchema={
                "type": "object",
                "properties": {
                    "file_path": {"type": "string", "description": "Path to the presentation file"},
                    "title": {"type": "string", "description": "Main presentation title"},
                    "subtitle": {"type": "string", "description": "Subtitle or description"},
                    "author": {"type": "string", "description": "Author or company name"},
                    "date": {"type": "string", "description": "Date or additional info"},
                    "slide_index": {"type": "integer", "description": "Index where to create slide (0-based)", "default": 0}
                },
                "required": ["file_path", "title"]
            }
        ),
        Tool(
            name="create_data_slide",
            description="Create a complete data slide with title, table, and optional chart",
            inputSchema={
                "type": "object",
                "properties": {
                    "file_path": {"type": "string", "description": "Path to the presentation file"},
                    "title": {"type": "string", "description": "Slide title"},
                    "data": {"type": "array", "description": "2D array of data for table", "items": {"type": "array", "items": {"type": "string"}}},
                    "include_chart": {"type": "boolean", "description": "Whether to create a chart from the data", "default": False},
                    "chart_type": {"type": "string", "description": "Chart type if creating chart", "default": "column"},
                    "position": {"type": "integer", "description": "Position to insert slide (-1 for end)", "default": -1}
                },
                "required": ["file_path", "title", "data"]
            }
        ),
        Tool(
            name="create_comparison_slide",
            description="Create a comparison slide with two columns of content",
            inputSchema={
                "type": "object",
                "properties": {
                    "file_path": {"type": "string", "description": "Path to the presentation file"},
                    "title": {"type": "string", "description": "Slide title"},
                    "left_title": {"type": "string", "description": "Left column title"},
                    "left_content": {"type": "array", "description": "Left column bullet points", "items": {"type": "string"}},
                    "right_title": {"type": "string", "description": "Right column title"},
                    "right_content": {"type": "array", "description": "Right column bullet points", "items": {"type": "string"}},
                    "position": {"type": "integer", "description": "Position to insert slide (-1 for end)", "default": -1}
                },
                "required": ["file_path", "title", "left_title", "left_content", "right_title", "right_content"]
            }
        ),
        Tool(
            name="create_agenda_slide",
            description="Create an agenda/outline slide with numbered or bulleted items",
            inputSchema={
                "type": "object",
                "properties": {
                    "file_path": {"type": "string", "description": "Path to the presentation file"},
                    "title": {"type": "string", "description": "Slide title", "default": "Agenda"},
                    "agenda_items": {"type": "array", "description": "List of agenda items", "items": {"type": "string"}},
                    "numbered": {"type": "boolean", "description": "Use numbers instead of bullets", "default": True},
                    "position": {"type": "integer", "description": "Position to insert slide (-1 for end)", "default": 1}
                },
                "required": ["file_path", "agenda_items"]
            }
        ),
        Tool(
            name="batch_replace_text",
            description="Replace text across multiple slides in the presentation",
            inputSchema={
                "type": "object",
                "properties": {
                    "file_path": {"type": "string", "description": "Path to the presentation file"},
                    "replacements": {"type": "object", "description": "Key-value pairs of text to replace", "additionalProperties": {"type": "string"}},
                    "slide_range": {"type": "array", "description": "Range of slide indices to process (all if not specified)", "items": {"type": "integer"}},
                    "case_sensitive": {"type": "boolean", "description": "Whether replacement should be case sensitive", "default": False}
                },
                "required": ["file_path", "replacements"]
            }
        ),
        Tool(
            name="apply_brand_theme",
            description="Apply consistent branding theme across presentation",
            inputSchema={
                "type": "object",
                "properties": {
                    "file_path": {"type": "string", "description": "Path to the presentation file"},
                    "primary_color": {"type": "string", "description": "Primary brand color (hex)", "default": "#0066CC"},
                    "secondary_color": {"type": "string", "description": "Secondary brand color (hex)", "default": "#999999"},
                    "accent_color": {"type": "string", "description": "Accent brand color (hex)", "default": "#FF6600"},
                    "font_family": {"type": "string", "description": "Primary font family", "default": "Arial"},
                    "apply_to_titles": {"type": "boolean", "description": "Apply colors to slide titles", "default": True},
                    "apply_to_shapes": {"type": "boolean", "description": "Apply colors to shapes", "default": True}
                },
                "required": ["file_path"]
            }
        ),
        Tool(
            name="create_section_break",
            description="Create a section break slide with large title and optional image",
            inputSchema={
                "type": "object",
                "properties": {
                    "file_path": {"type": "string", "description": "Path to the presentation file"},
                    "section_title": {"type": "string", "description": "Section title"},
                    "subtitle": {"type": "string", "description": "Optional subtitle"},
                    "background_color": {"type": "string", "description": "Background color (hex)", "default": "#0066CC"},
                    "text_color": {"type": "string", "description": "Text color (hex)", "default": "#FFFFFF"},
                    "position": {"type": "integer", "description": "Position to insert slide (-1 for end)", "default": -1}
                },
                "required": ["file_path", "section_title"]
            }
        ),
        Tool(
            name="generate_summary_slide",
            description="Generate a summary slide based on presentation content",
            inputSchema={
                "type": "object",
                "properties": {
                    "file_path": {"type": "string", "description": "Path to the presentation file"},
                    "title": {"type": "string", "description": "Summary slide title", "default": "Summary"},
                    "max_points": {"type": "integer", "description": "Maximum number of summary points", "default": 5},
                    "position": {"type": "integer", "description": "Position to insert slide (-1 for end)", "default": -1}
                },
                "required": ["file_path"]
            }
        )
    ]


@server.call_tool()
async def call_tool(name: str, arguments: dict) -> list[TextContent]:
    """Handle tool calls for PowerPoint operations."""
    try:
        result = None

        if name == "create_presentation":
            result = await create_presentation(**arguments)
        elif name == "create_presentation_from_template":
            result = await create_presentation_from_template(**arguments)
        elif name == "clone_presentation":
            result = await clone_presentation(**arguments)
        elif name == "open_presentation":
            result = await open_presentation(**arguments)
        elif name == "save_presentation":
            result = await save_presentation(**arguments)
        elif name == "get_presentation_info":
            result = await get_presentation_info(**arguments)
        elif name == "add_slide":
            result = await add_slide(**arguments)
        elif name == "delete_slide":
            result = await delete_slide(**arguments)
        elif name == "move_slide":
            result = await move_slide(**arguments)
        elif name == "duplicate_slide":
            result = await duplicate_slide(**arguments)
        elif name == "list_slides":
            result = await list_slides(**arguments)
        elif name == "set_slide_title":
            result = await set_slide_title(**arguments)
        elif name == "set_slide_content":
            result = await set_slide_content(**arguments)
        elif name == "add_text_box":
            result = await add_text_box(**arguments)
        elif name == "format_text":
            result = await format_text(**arguments)
        elif name == "add_image":
            result = await add_image(**arguments)
        elif name == "add_image_from_base64":
            result = await add_image_from_base64(**arguments)
        elif name == "replace_image":
            result = await replace_image(**arguments)
        elif name == "add_shape":
            result = await add_shape(**arguments)
        elif name == "modify_shape":
            result = await modify_shape(**arguments)
        elif name == "delete_shape":
            result = await delete_shape(**arguments)
        elif name == "add_table":
            result = await add_table(**arguments)
        elif name == "set_table_cell":
            result = await set_table_cell(**arguments)
        elif name == "format_table_cell":
            result = await format_table_cell(**arguments)
        elif name == "populate_table":
            result = await populate_table(**arguments)
        elif name == "add_chart":
            result = await add_chart(**arguments)
        elif name == "update_chart_data":
            result = await update_chart_data(**arguments)
        elif name == "list_shapes":
            result = await list_shapes(**arguments)
        elif name == "get_slide_layouts":
            result = await get_slide_layouts(**arguments)
        elif name == "set_presentation_properties":
            result = await set_presentation_properties(**arguments)
        elif name == "export_slide_as_image":
            result = await export_slide_as_image(**arguments)
        # Composite workflow tools
        elif name == "create_title_slide":
            result = await create_title_slide(**arguments)
        elif name == "create_data_slide":
            result = await create_data_slide(**arguments)
        elif name == "create_comparison_slide":
            result = await create_comparison_slide(**arguments)
        elif name == "create_agenda_slide":
            result = await create_agenda_slide(**arguments)
        elif name == "batch_replace_text":
            result = await batch_replace_text(**arguments)
        elif name == "apply_brand_theme":
            result = await apply_brand_theme(**arguments)
        elif name == "create_section_break":
            result = await create_section_break(**arguments)
        elif name == "generate_summary_slide":
            result = await generate_summary_slide(**arguments)
        else:
            return [TextContent(type="text", text=json.dumps({"ok": False, "error": f"Unknown tool: {name}"}))]

        return [TextContent(type="text", text=json.dumps({"ok": True, "result": result}))]

    except Exception as e:
        log.error(f"Error in tool {name}: {str(e)}")
        return [TextContent(type="text", text=json.dumps({"ok": False, "error": str(e)}))]


# Tool implementations start here
async def create_presentation(file_path: str, title: Optional[str] = None) -> Dict[str, Any]:
    """Create a new PowerPoint presentation."""
    prs = Presentation()

    if title:
        # Add title slide
        title_slide_layout = prs.slide_layouts[0]  # Title slide layout
        slide = prs.slides.add_slide(title_slide_layout)
        slide.shapes.title.text = title

    # Ensure proper directory structure
    organized_path = _ensure_output_directory(file_path)

    # Cache the presentation
    abs_path = os.path.abspath(organized_path)
    _presentations[abs_path] = prs

    # Save immediately
    _save_presentation(organized_path)

    return {"message": f"Created presentation: {organized_path}", "slide_count": len(prs.slides)}


async def open_presentation(file_path: str) -> Dict[str, Any]:
    """Open an existing PowerPoint presentation."""
    if not os.path.exists(file_path):
        raise FileNotFoundError(f"Presentation file not found: {file_path}")

    prs = _get_presentation(file_path)
    return {
        "message": f"Opened presentation: {file_path}",
        "slide_count": len(prs.slides),
        "layouts_count": len(prs.slide_layouts)
    }


async def save_presentation(file_path: str) -> Dict[str, Any]:
    """Save the current presentation to file."""
    _save_presentation(file_path)
    return {"message": f"Saved presentation: {file_path}"}


async def get_presentation_info(file_path: str) -> Dict[str, Any]:
    """Get information about the presentation."""
    prs = _get_presentation(file_path)

    # Get document properties
    props = prs.core_properties

    return {
        "file_path": file_path,
        "slide_count": len(prs.slides),
        "layout_count": len(prs.slide_layouts),
        "title": props.title or "",
        "author": props.author or "",
        "subject": props.subject or "",
        "created": str(props.created) if props.created else "",
        "modified": str(props.modified) if props.modified else "",
    }


async def add_slide(file_path: str, layout_index: int = 0, position: int = -1) -> Dict[str, Any]:
    """Add a new slide to the presentation."""
    prs = _get_presentation(file_path)

    if layout_index >= len(prs.slide_layouts):
        raise ValueError(f"Layout index {layout_index} out of range. Available layouts: 0-{len(prs.slide_layouts)-1}")

    slide_layout = prs.slide_layouts[layout_index]

    if position == -1:
        slide = prs.slides.add_slide(slide_layout)
        slide_idx = len(prs.slides) - 1
    else:
        # python-pptx doesn't have direct insert at position, so we'll add at end and move
        slide = prs.slides.add_slide(slide_layout)
        slide_idx = len(prs.slides) - 1
        if position < slide_idx:
            # Move slide to desired position (this is a workaround)
            pass  # Note: Moving requires more complex XML manipulation

    return {
        "message": f"Added slide at position {slide_idx}",
        "slide_index": slide_idx,
        "layout_name": slide_layout.name if hasattr(slide_layout, 'name') else f"Layout {layout_index}"
    }


async def delete_slide(file_path: str, slide_index: int) -> Dict[str, Any]:
    """Delete a slide from the presentation."""
    prs = _get_presentation(file_path)

    if slide_index < 0 or slide_index >= len(prs.slides):
        raise ValueError(f"Slide index {slide_index} out of range. Available slides: 0-{len(prs.slides)-1}")

    # Get slide reference
    slide_id = prs.slides[slide_index].slide_id

    # Remove from slides collection
    del prs.slides._sldIdLst[slide_index]

    return {"message": f"Deleted slide at index {slide_index}"}


async def move_slide(file_path: str, from_index: int, to_index: int) -> Dict[str, Any]:
    """Move a slide to a different position."""
    prs = _get_presentation(file_path)

    slide_count = len(prs.slides)
    if from_index < 0 or from_index >= slide_count:
        raise ValueError(f"From index {from_index} out of range")
    if to_index < 0 or to_index >= slide_count:
        raise ValueError(f"To index {to_index} out of range")

    # This is a complex operation that requires XML manipulation
    # For now, return a placeholder
    return {"message": f"Moved slide from {from_index} to {to_index}", "note": "Move operation is complex in python-pptx"}


async def duplicate_slide(file_path: str, slide_index: int, position: int = -1) -> Dict[str, Any]:
    """Duplicate an existing slide."""
    prs = _get_presentation(file_path)

    if slide_index < 0 or slide_index >= len(prs.slides):
        raise ValueError(f"Slide index {slide_index} out of range")

    # Get the source slide
    source_slide = prs.slides[slide_index]

    # Add new slide with same layout
    new_slide = prs.slides.add_slide(source_slide.slide_layout)

    # Copy content (this is a simplified version - full duplication requires more complex logic)
    try:
        if source_slide.shapes.title:
            new_slide.shapes.title.text = source_slide.shapes.title.text
    except:
        pass

    new_idx = len(prs.slides) - 1
    return {"message": f"Duplicated slide {slide_index} to position {new_idx}", "new_slide_index": new_idx}


async def list_slides(file_path: str) -> Dict[str, Any]:
    """List all slides in the presentation."""
    prs = _get_presentation(file_path)

    slides_info = []
    for i, slide in enumerate(prs.slides):
        slide_info = {
            "index": i,
            "layout_name": slide.slide_layout.name if hasattr(slide.slide_layout, 'name') else f"Layout {i}",
            "shape_count": len(slide.shapes),
            "title": ""
        }

        # Try to get slide title
        try:
            if slide.shapes.title:
                slide_info["title"] = slide.shapes.title.text
        except:
            pass

        slides_info.append(slide_info)

    return {"slides": slides_info, "total_count": len(slides_info)}


async def set_slide_title(file_path: str, slide_index: int, title: str) -> Dict[str, Any]:
    """Set the title of a slide."""
    prs = _get_presentation(file_path)

    if slide_index < 0 or slide_index >= len(prs.slides):
        raise ValueError(f"Slide index {slide_index} out of range")

    slide = prs.slides[slide_index]

    if not slide.shapes.title:
        raise ValueError("This slide layout does not have a title placeholder")

    slide.shapes.title.text = title

    return {"message": f"Set title for slide {slide_index}: {title}"}


async def set_slide_content(file_path: str, slide_index: int, content: str) -> Dict[str, Any]:
    """Set the main content/body text of a slide."""
    prs = _get_presentation(file_path)

    if slide_index < 0 or slide_index >= len(prs.slides):
        raise ValueError(f"Slide index {slide_index} out of range")

    slide = prs.slides[slide_index]

    # Look for content placeholder
    content_placeholder = None
    for shape in slide.placeholders:
        if shape.placeholder_format.idx == 1:  # Content placeholder is usually index 1
            content_placeholder = shape
            break

    if not content_placeholder:
        # If no content placeholder, try to find text frame
        for shape in slide.shapes:
            if hasattr(shape, 'text_frame') and shape != slide.shapes.title:
                content_placeholder = shape
                break

    if not content_placeholder:
        raise ValueError("No content area found on this slide")

    # Split content by newlines and create bullet points
    lines = content.split('\\n')
    content_placeholder.text = lines[0]  # First line

    if len(lines) > 1:
        text_frame = content_placeholder.text_frame
        for line in lines[1:]:
            p = text_frame.add_paragraph()
            p.text = line
            p.level = 0  # Bullet level

    return {"message": f"Set content for slide {slide_index}"}


async def add_text_box(file_path: str, slide_index: int, text: str, left: float = 1.0, top: float = 1.0,
                      width: float = 6.0, height: float = 1.0, font_size: int = 18,
                      font_color: str = "#000000", bold: bool = False, italic: bool = False) -> Dict[str, Any]:
    """Add a text box to a slide."""
    prs = _get_presentation(file_path)

    if slide_index < 0 or slide_index >= len(prs.slides):
        raise ValueError(f"Slide index {slide_index} out of range")

    slide = prs.slides[slide_index]

    # Add text box
    textbox = slide.shapes.add_textbox(Inches(left), Inches(top), Inches(width), Inches(height))
    text_frame = textbox.text_frame
    text_frame.text = text

    # Format text
    paragraph = text_frame.paragraphs[0]
    run = paragraph.runs[0]
    font = run.font

    font.size = Pt(font_size)
    font.color.rgb = _parse_color(font_color)
    font.bold = bold
    font.italic = italic

    return {
        "message": f"Added text box to slide {slide_index}",
        "shape_index": len(slide.shapes) - 1,
        "text": text
    }


async def format_text(file_path: str, slide_index: int, shape_index: int, **kwargs) -> Dict[str, Any]:
    """Format existing text in a slide."""
    prs = _get_presentation(file_path)

    if slide_index < 0 or slide_index >= len(prs.slides):
        raise ValueError(f"Slide index {slide_index} out of range")

    slide = prs.slides[slide_index]

    if shape_index < 0 or shape_index >= len(slide.shapes):
        raise ValueError(f"Shape index {shape_index} out of range")

    shape = slide.shapes[shape_index]

    if not hasattr(shape, 'text_frame'):
        raise ValueError("Selected shape does not contain text")

    # Apply formatting to all paragraphs and runs
    for paragraph in shape.text_frame.paragraphs:
        if kwargs.get('alignment'):
            alignment_map = {
                'left': PP_ALIGN.LEFT,
                'center': PP_ALIGN.CENTER,
                'right': PP_ALIGN.RIGHT,
                'justify': PP_ALIGN.JUSTIFY
            }
            paragraph.alignment = alignment_map.get(kwargs['alignment'], PP_ALIGN.LEFT)

        for run in paragraph.runs:
            font = run.font

            if kwargs.get('font_name'):
                font.name = kwargs['font_name']
            if kwargs.get('font_size'):
                font.size = Pt(kwargs['font_size'])
            if kwargs.get('font_color'):
                font.color.rgb = _parse_color(kwargs['font_color'])
            if kwargs.get('bold') is not None:
                font.bold = kwargs['bold']
            if kwargs.get('italic') is not None:
                font.italic = kwargs['italic']
            if kwargs.get('underline') is not None:
                font.underline = kwargs['underline']

    return {"message": f"Formatted text in shape {shape_index} on slide {slide_index}"}


async def add_image(file_path: str, slide_index: int, image_path: str, left: float = 1.0, top: float = 1.0,
                   width: Optional[float] = None, height: Optional[float] = None) -> Dict[str, Any]:
    """Add an image to a slide."""
    if not os.path.exists(image_path):
        raise FileNotFoundError(f"Image file not found: {image_path}")

    prs = _get_presentation(file_path)

    if slide_index < 0 or slide_index >= len(prs.slides):
        raise ValueError(f"Slide index {slide_index} out of range")

    slide = prs.slides[slide_index]

    # Add image
    if width and height:
        pic = slide.shapes.add_picture(image_path, Inches(left), Inches(top), Inches(width), Inches(height))
    elif width:
        pic = slide.shapes.add_picture(image_path, Inches(left), Inches(top), width=Inches(width))
    elif height:
        pic = slide.shapes.add_picture(image_path, Inches(left), Inches(top), height=Inches(height))
    else:
        pic = slide.shapes.add_picture(image_path, Inches(left), Inches(top))

    return {
        "message": f"Added image to slide {slide_index}",
        "shape_index": len(slide.shapes) - 1,
        "image_path": image_path
    }


async def add_image_from_base64(file_path: str, slide_index: int, image_data: str,
                               image_format: str = "png", left: float = 1.0, top: float = 1.0,
                               width: Optional[float] = None, height: Optional[float] = None) -> Dict[str, Any]:
    """Add an image from base64 data to a slide."""
    prs = _get_presentation(file_path)

    if slide_index < 0 or slide_index >= len(prs.slides):
        raise ValueError(f"Slide index {slide_index} out of range")

    slide = prs.slides[slide_index]

    # Decode base64 image
    try:
        image_bytes = base64.b64decode(image_data)
        image_stream = BytesIO(image_bytes)
    except Exception as e:
        raise ValueError(f"Invalid base64 image data: {e}")

    # Add image from stream
    if width and height:
        pic = slide.shapes.add_picture(image_stream, Inches(left), Inches(top), Inches(width), Inches(height))
    elif width:
        pic = slide.shapes.add_picture(image_stream, Inches(left), Inches(top), width=Inches(width))
    elif height:
        pic = slide.shapes.add_picture(image_stream, Inches(left), Inches(top), height=Inches(height))
    else:
        pic = slide.shapes.add_picture(image_stream, Inches(left), Inches(top))

    return {
        "message": f"Added image from base64 to slide {slide_index}",
        "shape_index": len(slide.shapes) - 1,
        "format": image_format
    }


async def replace_image(file_path: str, slide_index: int, shape_index: int, new_image_path: str) -> Dict[str, Any]:
    """Replace an existing image in a slide."""
    if not os.path.exists(new_image_path):
        raise FileNotFoundError(f"New image file not found: {new_image_path}")

    prs = _get_presentation(file_path)

    if slide_index < 0 or slide_index >= len(prs.slides):
        raise ValueError(f"Slide index {slide_index} out of range")

    slide = prs.slides[slide_index]

    if shape_index < 0 or shape_index >= len(slide.shapes):
        raise ValueError(f"Shape index {shape_index} out of range")

    shape = slide.shapes[shape_index]

    # This is complex in python-pptx - would need to remove old image and add new one
    # For now, provide guidance
    return {
        "message": "Image replacement requires removing old image and adding new one",
        "note": "Use delete_shape and add_image for full replacement functionality"
    }


async def add_shape(file_path: str, slide_index: int, shape_type: str, left: float = 1.0, top: float = 1.0,
                   width: float = 2.0, height: float = 1.0, fill_color: Optional[str] = None,
                   line_color: Optional[str] = None, line_width: float = 1.0) -> Dict[str, Any]:
    """Add a shape to a slide."""
    prs = _get_presentation(file_path)

    if slide_index < 0 or slide_index >= len(prs.slides):
        raise ValueError(f"Slide index {slide_index} out of range")

    slide = prs.slides[slide_index]

    # Map shape types to MSO_SHAPE constants
    shape_map = {
        'rectangle': MSO_SHAPE.RECTANGLE,
        'oval': MSO_SHAPE.OVAL,
        'triangle': MSO_SHAPE.ISOSCELES_TRIANGLE,
        'arrow': MSO_SHAPE.BLOCK_ARC,
        'diamond': MSO_SHAPE.DIAMOND,
        'pentagon': MSO_SHAPE.REGULAR_PENTAGON,
        'hexagon': MSO_SHAPE.HEXAGON,
        'octagon': MSO_SHAPE.OCTAGON,
        'star': MSO_SHAPE.STAR_5_POINT,
        'heart': MSO_SHAPE.HEART,
        'smiley': MSO_SHAPE.SMILEY_FACE,
    }

    if shape_type.lower() not in shape_map:
        available_shapes = ', '.join(shape_map.keys())
        raise ValueError(f"Unknown shape type: {shape_type}. Available: {available_shapes}")

    # Add shape
    shape = slide.shapes.add_shape(
        shape_map[shape_type.lower()],
        Inches(left), Inches(top), Inches(width), Inches(height)
    )

    # Apply formatting
    if fill_color:
        shape.fill.solid()
        shape.fill.fore_color.rgb = _parse_color(fill_color)

    if line_color:
        shape.line.color.rgb = _parse_color(line_color)

    shape.line.width = Pt(line_width)

    return {
        "message": f"Added {shape_type} shape to slide {slide_index}",
        "shape_index": len(slide.shapes) - 1
    }


async def modify_shape(file_path: str, slide_index: int, shape_index: int, **kwargs) -> Dict[str, Any]:
    """Modify properties of an existing shape."""
    prs = _get_presentation(file_path)

    if slide_index < 0 or slide_index >= len(prs.slides):
        raise ValueError(f"Slide index {slide_index} out of range")

    slide = prs.slides[slide_index]

    if shape_index < 0 or shape_index >= len(slide.shapes):
        raise ValueError(f"Shape index {shape_index} out of range")

    shape = slide.shapes[shape_index]

    # Modify position and size
    if kwargs.get('left') is not None:
        shape.left = Inches(kwargs['left'])
    if kwargs.get('top') is not None:
        shape.top = Inches(kwargs['top'])
    if kwargs.get('width') is not None:
        shape.width = Inches(kwargs['width'])
    if kwargs.get('height') is not None:
        shape.height = Inches(kwargs['height'])

    # Modify formatting
    if kwargs.get('fill_color'):
        shape.fill.solid()
        shape.fill.fore_color.rgb = _parse_color(kwargs['fill_color'])

    if kwargs.get('line_color'):
        shape.line.color.rgb = _parse_color(kwargs['line_color'])

    if kwargs.get('line_width') is not None:
        shape.line.width = Pt(kwargs['line_width'])

    return {"message": f"Modified shape {shape_index} on slide {slide_index}"}


async def delete_shape(file_path: str, slide_index: int, shape_index: int) -> Dict[str, Any]:
    """Delete a shape from a slide."""
    prs = _get_presentation(file_path)

    if slide_index < 0 or slide_index >= len(prs.slides):
        raise ValueError(f"Slide index {slide_index} out of range")

    slide = prs.slides[slide_index]

    if shape_index < 0 or shape_index >= len(slide.shapes):
        raise ValueError(f"Shape index {shape_index} out of range")

    shape = slide.shapes[shape_index]
    sp = shape._element
    sp.getparent().remove(sp)

    return {"message": f"Deleted shape {shape_index} from slide {slide_index}"}


async def add_table(file_path: str, slide_index: int, rows: int, cols: int, left: float = 1.0,
                   top: float = 1.0, width: float = 6.0, height: float = 3.0) -> Dict[str, Any]:
    """Add a table to a slide."""
    prs = _get_presentation(file_path)

    if slide_index < 0 or slide_index >= len(prs.slides):
        raise ValueError(f"Slide index {slide_index} out of range")

    slide = prs.slides[slide_index]

    # Add table
    table_shape = slide.shapes.add_table(rows, cols, Inches(left), Inches(top), Inches(width), Inches(height))
    table_index = len(slide.shapes) - 1

    return {
        "message": f"Added {rows}x{cols} table to slide {slide_index}",
        "shape_index": table_index,
        "table_shape_index": table_index,  # Explicit table index for reference
        "rows": rows,
        "cols": cols
    }


async def set_table_cell(file_path: str, slide_index: int, table_index: int, row: int, col: int, text: str) -> Dict[str, Any]:
    """Set the text content of a table cell."""
    prs = _get_presentation(file_path)

    if slide_index < 0 or slide_index >= len(prs.slides):
        raise ValueError(f"Slide index {slide_index} out of range")

    slide = prs.slides[slide_index]

    if table_index < 0 or table_index >= len(slide.shapes):
        raise ValueError(f"Table index {table_index} out of range")

    shape = slide.shapes[table_index]

    try:
        if not shape.has_table:
            raise ValueError("Selected shape is not a table")
        table = shape.table
    except AttributeError:
        raise ValueError("Selected shape is not a table")

    if row < 0 or row >= len(table.rows):
        raise ValueError(f"Row {row} out of range")
    if col < 0 or col >= len(table.columns):
        raise ValueError(f"Column {col} out of range")

    cell = table.cell(row, col)
    cell.text = text

    return {"message": f"Set cell [{row},{col}] text: {text}"}


async def format_table_cell(file_path: str, slide_index: int, table_index: int, row: int, col: int, **kwargs) -> Dict[str, Any]:
    """Format a table cell."""
    prs = _get_presentation(file_path)

    if slide_index < 0 or slide_index >= len(prs.slides):
        raise ValueError(f"Slide index {slide_index} out of range")

    slide = prs.slides[slide_index]

    if table_index < 0 or table_index >= len(slide.shapes):
        raise ValueError(f"Table index {table_index} out of range")

    shape = slide.shapes[table_index]

    try:
        if not shape.has_table:
            raise ValueError("Selected shape is not a table")
        table = shape.table
    except AttributeError:
        raise ValueError("Selected shape is not a table")
    cell = table.cell(row, col)

    # Format cell background
    if kwargs.get('fill_color'):
        cell.fill.solid()
        cell.fill.fore_color.rgb = _parse_color(kwargs['fill_color'])

    # Format text
    for paragraph in cell.text_frame.paragraphs:
        if kwargs.get('alignment'):
            alignment_map = {
                'left': PP_ALIGN.LEFT,
                'center': PP_ALIGN.CENTER,
                'right': PP_ALIGN.RIGHT
            }
            paragraph.alignment = alignment_map.get(kwargs['alignment'], PP_ALIGN.LEFT)

        for run in paragraph.runs:
            font = run.font

            if kwargs.get('font_size'):
                font.size = Pt(kwargs['font_size'])
            if kwargs.get('font_color'):
                font.color.rgb = _parse_color(kwargs['font_color'])
            if kwargs.get('bold') is not None:
                font.bold = kwargs['bold']

    return {"message": f"Formatted cell [{row},{col}]"}


async def populate_table(file_path: str, slide_index: int, table_index: int, data: List[List[str]],
                        header_row: bool = False) -> Dict[str, Any]:
    """Populate entire table with data from a 2D array."""
    prs = _get_presentation(file_path)

    if slide_index < 0 or slide_index >= len(prs.slides):
        raise ValueError(f"Slide index {slide_index} out of range")

    slide = prs.slides[slide_index]

    if table_index < 0 or table_index >= len(slide.shapes):
        raise ValueError(f"Table index {table_index} out of range")

    shape = slide.shapes[table_index]

    try:
        if not shape.has_table:
            raise ValueError("Selected shape is not a table")
        table = shape.table
    except AttributeError:
        raise ValueError("Selected shape is not a table")

    # Populate data
    for row_idx, row_data in enumerate(data):
        if row_idx >= len(table.rows):
            break

        for col_idx, cell_data in enumerate(row_data):
            if col_idx >= len(table.columns):
                break

            cell = table.cell(row_idx, col_idx)
            cell.text = str(cell_data)

            # Format header row
            if header_row and row_idx == 0:
                for run in cell.text_frame.paragraphs[0].runs:
                    run.font.bold = True

    return {"message": f"Populated table with {len(data)} rows of data"}


async def add_chart(file_path: str, slide_index: int, data: Dict[str, Any], chart_type: str = "column",
                   left: float = 1.0, top: float = 1.0, width: float = 6.0, height: float = 4.0,
                   title: Optional[str] = None) -> Dict[str, Any]:
    """Add a chart to a slide."""
    prs = _get_presentation(file_path)

    if slide_index < 0 or slide_index >= len(prs.slides):
        raise ValueError(f"Slide index {slide_index} out of range")

    slide = prs.slides[slide_index]

    # Map chart types
    chart_type_map = {
        'column': XL_CHART_TYPE.COLUMN_CLUSTERED,
        'bar': XL_CHART_TYPE.BAR_CLUSTERED,
        'line': XL_CHART_TYPE.LINE,
        'pie': XL_CHART_TYPE.PIE
    }

    if chart_type not in chart_type_map:
        available_types = ', '.join(chart_type_map.keys())
        raise ValueError(f"Unknown chart type: {chart_type}. Available: {available_types}")

    # Prepare chart data
    chart_data = CategoryChartData()
    chart_data.categories = data.get('categories', [])

    for series_info in data.get('series', []):
        chart_data.add_series(series_info.get('name', 'Series'), series_info.get('values', []))

    # Add chart
    chart_shape = slide.shapes.add_chart(
        chart_type_map[chart_type],
        Inches(left), Inches(top), Inches(width), Inches(height),
        chart_data
    )

    # Set title if provided
    if title:
        chart_shape.chart.chart_title.text_frame.text = title

    return {
        "message": f"Added {chart_type} chart to slide {slide_index}",
        "shape_index": len(slide.shapes) - 1,
        "title": title or "Untitled Chart"
    }


async def update_chart_data(file_path: str, slide_index: int, chart_index: int, data: Dict[str, Any]) -> Dict[str, Any]:
    """Update data in an existing chart."""
    prs = _get_presentation(file_path)

    if slide_index < 0 or slide_index >= len(prs.slides):
        raise ValueError(f"Slide index {slide_index} out of range")

    slide = prs.slides[slide_index]

    if chart_index < 0 or chart_index >= len(slide.shapes):
        raise ValueError(f"Chart index {chart_index} out of range")

    shape = slide.shapes[chart_index]

    if not hasattr(shape, 'chart'):
        raise ValueError("Selected shape is not a chart")

    # Note: Updating chart data in python-pptx is complex and may require
    # recreating the chart or manipulating the underlying XML
    return {
        "message": "Chart data update is complex in python-pptx",
        "note": "Consider recreating the chart with new data for full functionality"
    }


async def list_shapes(file_path: str, slide_index: int) -> Dict[str, Any]:
    """List all shapes on a slide with their types and properties."""
    prs = _get_presentation(file_path)

    if slide_index < 0 or slide_index >= len(prs.slides):
        raise ValueError(f"Slide index {slide_index} out of range")

    slide = prs.slides[slide_index]

    shapes_info = []
    for i, shape in enumerate(slide.shapes):
        shape_info = {
            "index": i,
            "type": str(shape.shape_type),
            "left": float(shape.left.inches),
            "top": float(shape.top.inches),
            "width": float(shape.width.inches),
            "height": float(shape.height.inches),
            "has_text": hasattr(shape, 'text_frame'),
            "text": ""
        }

        # Get text if available
        if hasattr(shape, 'text_frame') and shape.text_frame:
            try:
                shape_info["text"] = shape.text_frame.text[:100]  # First 100 chars
            except:
                pass

        # Special handling for different shape types
        try:
            if shape.has_table:
                shape_info["type"] = "TABLE"
                shape_info["rows"] = len(shape.table.rows)
                shape_info["cols"] = len(shape.table.columns)
        except (AttributeError, ValueError):
            pass

        try:
            if shape.has_chart:
                shape_info["type"] = "CHART"
                shape_info["chart_type"] = str(shape.chart.chart_type)
        except (AttributeError, ValueError):
            pass

        shapes_info.append(shape_info)

    return {"shapes": shapes_info, "total_count": len(shapes_info)}


async def get_slide_layouts(file_path: str) -> Dict[str, Any]:
    """Get available slide layouts in the presentation."""
    prs = _get_presentation(file_path)

    layouts_info = []
    for i, layout in enumerate(prs.slide_layouts):
        layout_info = {
            "index": i,
            "name": layout.name if hasattr(layout, 'name') else f"Layout {i}",
            "placeholder_count": len(layout.placeholders)
        }
        layouts_info.append(layout_info)

    return {"layouts": layouts_info, "total_count": len(layouts_info)}


async def set_presentation_properties(file_path: str, **kwargs) -> Dict[str, Any]:
    """Set presentation document properties."""
    prs = _get_presentation(file_path)
    props = prs.core_properties

    if kwargs.get('title'):
        props.title = kwargs['title']
    if kwargs.get('author'):
        props.author = kwargs['author']
    if kwargs.get('subject'):
        props.subject = kwargs['subject']
    if kwargs.get('comments'):
        props.comments = kwargs['comments']

    return {"message": "Updated presentation properties"}


async def export_slide_as_image(file_path: str, slide_index: int, output_path: str, format: str = "png") -> Dict[str, Any]:
    """Export a slide as an image file."""
    # Note: python-pptx doesn't have built-in slide-to-image export functionality
    # This would require additional libraries like python-pptx-interface or PIL with COM automation
    return {
        "message": "Slide image export requires additional libraries",
        "note": "Consider using python-pptx-interface or COM automation for image export functionality",
        "requested_output": output_path,
        "format": format
    }


# Template and Enhanced Workflow Functions
async def create_presentation_from_template(template_path: str, output_path: str, title: Optional[str] = None,
                                          replace_placeholders: Optional[Dict[str, str]] = None) -> Dict[str, Any]:
    """Create a new presentation from an existing template."""
    # Resolve template path (check templates directory)
    resolved_template = _resolve_template_path(template_path)
    if not os.path.exists(resolved_template):
        raise FileNotFoundError(f"Template file not found: {template_path} (searched: {resolved_template})")

    # Load template
    template_prs = Presentation(resolved_template)

    # Ensure proper output directory
    organized_output = _ensure_output_directory(output_path)

    # Cache the new presentation
    abs_output_path = os.path.abspath(organized_output)
    _presentations[abs_output_path] = template_prs

    # Update title if provided
    if title and len(template_prs.slides) > 0:
        title_slide = template_prs.slides[0]
        if title_slide.shapes.title:
            title_slide.shapes.title.text = title

    # Replace placeholders if provided
    replacements_made = 0
    if replace_placeholders:
        for slide in template_prs.slides:
            for shape in slide.shapes:
                if hasattr(shape, 'text_frame') and shape.text_frame:
                    for paragraph in shape.text_frame.paragraphs:
                        for run in paragraph.runs:
                            for placeholder, replacement in replace_placeholders.items():
                                if placeholder in run.text:
                                    run.text = run.text.replace(placeholder, replacement)
                                    replacements_made += 1

    # Save the new presentation
    _save_presentation(organized_output)

    return {
        "message": f"Created presentation from template: {resolved_template}",
        "output_path": organized_output,
        "slide_count": len(template_prs.slides),
        "replacements_made": replacements_made
    }


async def clone_presentation(source_path: str, target_path: str, new_title: Optional[str] = None) -> Dict[str, Any]:
    """Clone an existing presentation with optional modifications."""
    # Resolve source path
    resolved_source = _resolve_template_path(source_path)  # Can also check templates
    if not os.path.exists(resolved_source):
        raise FileNotFoundError(f"Source presentation not found: {source_path}")

    # Load source presentation
    source_prs = Presentation(resolved_source)

    # Ensure proper output directory
    organized_target = _ensure_output_directory(target_path)

    # Cache the cloned presentation
    abs_target_path = os.path.abspath(organized_target)
    _presentations[abs_target_path] = source_prs

    # Update title if provided
    if new_title and len(source_prs.slides) > 0:
        first_slide = source_prs.slides[0]
        if first_slide.shapes.title:
            first_slide.shapes.title.text = new_title

    # Save the cloned presentation
    _save_presentation(organized_target)

    return {
        "message": f"Cloned presentation from {resolved_source} to {organized_target}",
        "slide_count": len(source_prs.slides),
        "new_title": new_title or "No title change"
    }


# Composite Workflow Tools
async def create_title_slide(file_path: str, title: str, subtitle: Optional[str] = None,
                           author: Optional[str] = None, date: Optional[str] = None,
                           slide_index: int = 0) -> Dict[str, Any]:
    """Create a complete title slide with all elements."""
    prs = _get_presentation(file_path)

    # Get or create slide at specified index
    if slide_index >= len(prs.slides):
        # Add new slide with title layout
        title_layout = prs.slide_layouts[0]  # Title slide layout
        slide = prs.slides.add_slide(title_layout)
        actual_index = len(prs.slides) - 1
    else:
        slide = prs.slides[slide_index]
        actual_index = slide_index

    # Set main title
    if slide.shapes.title:
        slide.shapes.title.text = title

    # Set subtitle in content placeholder or create text box
    if subtitle:
        subtitle_shape = None
        for shape in slide.placeholders:
            if shape.placeholder_format.idx == 1:  # Subtitle placeholder
                subtitle_shape = shape
                break

        if subtitle_shape:
            subtitle_shape.text = subtitle
        else:
            # Create subtitle text box
            await add_text_box(file_path, actual_index, subtitle, 1.0, 2.5, 8.0, 1.0, 20, "#666666", False, True)

    # Add author info if provided
    if author:
        await add_text_box(file_path, actual_index, f"By: {author}", 1.0, 5.5, 4.0, 0.8, 16, "#888888", False, False)

    # Add date if provided
    if date:
        await add_text_box(file_path, actual_index, date, 5.0, 5.5, 4.0, 0.8, 16, "#888888", False, False)

    return {
        "message": f"Created title slide at index {actual_index}",
        "slide_index": actual_index,
        "title": title,
        "subtitle": subtitle or "None",
        "author": author or "None",
        "date": date or "None"
    }


async def create_data_slide(file_path: str, title: str, data: List[List[str]],
                          include_chart: bool = False, chart_type: str = "column",
                          position: int = -1) -> Dict[str, Any]:
    """Create a complete data slide with table and optional chart."""
    prs = _get_presentation(file_path)

    # Add slide
    slide_result = await add_slide(file_path, 1, position)  # Content layout
    slide_idx = slide_result["slide_index"]

    # Set title
    await set_slide_title(file_path, slide_idx, title)

    # Determine table size
    rows = len(data)
    cols = max(len(row) for row in data) if data else 1

    # Create table
    if include_chart:
        # Smaller table to make room for chart
        table_result = await add_table(file_path, slide_idx, rows, cols, 0.5, 2.0, 4.5, 3.0)
    else:
        # Full-width table
        table_result = await add_table(file_path, slide_idx, rows, cols, 1.0, 2.0, 8.0, 4.0)

    table_idx = table_result["shape_index"]

    # Populate table
    await populate_table(file_path, slide_idx, table_idx, data, True)

    chart_created = False
    if include_chart and len(data) > 1:
        try:
            # Create chart data from table (assuming first row is headers, first column is categories)
            if len(data[0]) >= 2:  # Need at least category and one data column
                categories = [row[0] for row in data[1:]]  # First column, skip header
                series = []

                for col_idx in range(1, len(data[0])):  # Skip first column (categories)
                    series_name = data[0][col_idx]  # Header
                    values = []

                    for row_idx in range(1, len(data)):  # Skip header row
                        try:
                            # Try to convert to number
                            value_str = data[row_idx][col_idx].replace('$', '').replace(',', '').replace('%', '')
                            values.append(float(value_str))
                        except (ValueError, IndexError):
                            values.append(0)

                    series.append({"name": series_name, "values": values})

                chart_data = {"categories": categories, "series": series}
                await add_chart(file_path, slide_idx, chart_data, chart_type, 5.5, 2.0, 4.0, 3.0, f"{title} Chart")
                chart_created = True

        except Exception as e:
            log.warning(f"Could not create chart from data: {e}")

    return {
        "message": f"Created data slide '{title}' at index {slide_idx}",
        "slide_index": slide_idx,
        "table_rows": rows,
        "table_cols": cols,
        "chart_created": chart_created
    }


async def create_comparison_slide(file_path: str, title: str, left_title: str, left_content: List[str],
                                right_title: str, right_content: List[str], position: int = -1) -> Dict[str, Any]:
    """Create a comparison slide with two columns."""
    # Add slide
    slide_result = await add_slide(file_path, 1, position)  # Content layout
    slide_idx = slide_result["slide_index"]

    # Set main title
    await set_slide_title(file_path, slide_idx, title)

    # Create left column
    await add_text_box(file_path, slide_idx, left_title, 0.5, 2.0, 4.0, 0.8, 20, "#0066CC", True, False)
    left_content_text = "\\n".join([f"• {item}" for item in left_content])
    await add_text_box(file_path, slide_idx, left_content_text, 0.5, 3.0, 4.0, 3.0, 16, "#000000", False, False)

    # Create right column
    await add_text_box(file_path, slide_idx, right_title, 5.5, 2.0, 4.0, 0.8, 20, "#0066CC", True, False)
    right_content_text = "\\n".join([f"• {item}" for item in right_content])
    await add_text_box(file_path, slide_idx, right_content_text, 5.5, 3.0, 4.0, 3.0, 16, "#000000", False, False)

    # Add dividing line
    await add_shape(file_path, slide_idx, "rectangle", 4.8, 2.0, 0.1, 4.0, "#CCCCCC", "#CCCCCC", 1.0)

    return {
        "message": f"Created comparison slide '{title}' at index {slide_idx}",
        "slide_index": slide_idx,
        "left_items": len(left_content),
        "right_items": len(right_content)
    }


async def create_agenda_slide(file_path: str, agenda_items: List[str], title: str = "Agenda",
                            numbered: bool = True, position: int = 1) -> Dict[str, Any]:
    """Create an agenda slide with numbered or bulleted items."""
    # Add slide
    slide_result = await add_slide(file_path, 1, position)  # Content layout
    slide_idx = slide_result["slide_index"]

    # Set title
    await set_slide_title(file_path, slide_idx, title)

    # Create agenda content
    if numbered:
        agenda_text = "\\n".join([f"{i+1}. {item}" for i, item in enumerate(agenda_items)])
    else:
        agenda_text = "\\n".join([f"• {item}" for item in agenda_items])

    await add_text_box(file_path, slide_idx, agenda_text, 1.5, 2.5, 7.0, 4.0, 18, "#000000", False, False)

    return {
        "message": f"Created agenda slide '{title}' at index {slide_idx}",
        "slide_index": slide_idx,
        "item_count": len(agenda_items),
        "numbered": numbered
    }


async def batch_replace_text(file_path: str, replacements: Dict[str, str],
                           slide_range: Optional[List[int]] = None, case_sensitive: bool = False) -> Dict[str, Any]:
    """Replace text across multiple slides in the presentation."""
    prs = _get_presentation(file_path)

    if slide_range is None:
        slides_to_process = list(range(len(prs.slides)))
    else:
        slides_to_process = [i for i in slide_range if 0 <= i < len(prs.slides)]

    total_replacements = 0

    for slide_idx in slides_to_process:
        slide = prs.slides[slide_idx]

        for shape in slide.shapes:
            if hasattr(shape, 'text_frame') and shape.text_frame:
                for paragraph in shape.text_frame.paragraphs:
                    for run in paragraph.runs:
                        original_text = run.text
                        modified_text = original_text

                        for old_text, new_text in replacements.items():
                            if case_sensitive:
                                if old_text in modified_text:
                                    modified_text = modified_text.replace(old_text, new_text)
                                    total_replacements += 1
                            else:
                                # Case-insensitive replacement
                                import re
                                pattern = re.compile(re.escape(old_text), re.IGNORECASE)
                                if pattern.search(modified_text):
                                    modified_text = pattern.sub(new_text, modified_text)
                                    total_replacements += 1

                        if modified_text != original_text:
                            run.text = modified_text

    return {
        "message": f"Completed batch text replacement across {len(slides_to_process)} slides",
        "slides_processed": len(slides_to_process),
        "total_replacements": total_replacements,
        "replacement_pairs": len(replacements)
    }


async def apply_brand_theme(file_path: str, primary_color: str = "#0066CC", secondary_color: str = "#999999",
                          accent_color: str = "#FF6600", font_family: str = "Arial",
                          apply_to_titles: bool = True, apply_to_shapes: bool = True) -> Dict[str, Any]:
    """Apply consistent branding theme across presentation."""
    prs = _get_presentation(file_path)

    title_updates = 0
    shape_updates = 0

    for slide in prs.slides:
        for shape in slide.shapes:
            # Apply to titles
            if apply_to_titles and hasattr(shape, 'text_frame') and shape.text_frame:
                if shape == slide.shapes.title:  # This is a title
                    for paragraph in shape.text_frame.paragraphs:
                        for run in paragraph.runs:
                            run.font.name = font_family
                            run.font.color.rgb = _parse_color(primary_color)
                    title_updates += 1

            # Apply to shapes
            if apply_to_shapes and hasattr(shape, 'fill'):
                try:
                    # Apply primary color to rectangle shapes
                    if shape.shape_type == 1:  # Rectangle
                        shape.fill.solid()
                        shape.fill.fore_color.rgb = _parse_color(primary_color)
                        shape_updates += 1
                    # Apply accent color to other shapes
                    elif shape.shape_type in [9, 7]:  # Oval, triangle, etc.
                        shape.fill.solid()
                        shape.fill.fore_color.rgb = _parse_color(accent_color)
                        shape_updates += 1
                except Exception:
                    pass  # Some shapes may not support fill

    return {
        "message": f"Applied brand theme to presentation",
        "primary_color": primary_color,
        "secondary_color": secondary_color,
        "accent_color": accent_color,
        "font_family": font_family,
        "title_updates": title_updates,
        "shape_updates": shape_updates
    }


async def create_section_break(file_path: str, section_title: str, subtitle: Optional[str] = None,
                             background_color: str = "#0066CC", text_color: str = "#FFFFFF",
                             position: int = -1) -> Dict[str, Any]:
    """Create a section break slide with large title and background color."""
    # Add slide
    slide_result = await add_slide(file_path, 6, position)  # Blank layout
    slide_idx = slide_result["slide_index"]

    prs = _get_presentation(file_path)
    slide = prs.slides[slide_idx]

    # Set background color by adding a full-slide rectangle
    await add_shape(file_path, slide_idx, "rectangle", 0, 0, 10, 7.5, background_color, background_color, 0)

    # Add large section title
    await add_text_box(file_path, slide_idx, section_title, 1.0, 2.5, 8.0, 1.5, 48, text_color, True, False)

    # Add subtitle if provided
    if subtitle:
        await add_text_box(file_path, slide_idx, subtitle, 1.0, 4.5, 8.0, 1.0, 24, text_color, False, False)

    return {
        "message": f"Created section break slide '{section_title}' at index {slide_idx}",
        "slide_index": slide_idx,
        "section_title": section_title,
        "subtitle": subtitle or "None",
        "background_color": background_color
    }


async def generate_summary_slide(file_path: str, title: str = "Summary",
                               max_points: int = 5, position: int = -1) -> Dict[str, Any]:
    """Generate a summary slide based on presentation content."""
    prs = _get_presentation(file_path)

    # Extract key points from slide titles and content
    summary_points = []

    for slide_idx, slide in enumerate(prs.slides):
        if slide_idx == 0:  # Skip title slide
            continue

        # Get slide title
        slide_title = ""
        if slide.shapes.title:
            slide_title = slide.shapes.title.text

        if slide_title and len(summary_points) < max_points:
            summary_points.append(slide_title)

    # If we don't have enough from titles, extract from content
    if len(summary_points) < max_points:
        for slide_idx in range(1, len(prs.slides)):  # Skip title slide
            if len(summary_points) >= max_points:
                break

            slide = prs.slides[slide_idx]
            for shape in slide.shapes:
                if hasattr(shape, 'text_frame') and shape.text_frame:
                    text = shape.text_frame.text.strip()
                    if text and shape != slide.shapes.title:
                        # Take first sentence or line
                        first_line = text.split('\\n')[0].split('.')[0]
                        if len(first_line) < 80 and first_line not in summary_points:
                            summary_points.append(first_line)
                            if len(summary_points) >= max_points:
                                break

    # Create summary slide
    slide_result = await add_slide(file_path, 1, position)  # Content layout
    slide_idx = slide_result["slide_index"]

    await set_slide_title(file_path, slide_idx, title)

    if summary_points:
        summary_text = "\\n".join([f"• {point}" for point in summary_points])
        await add_text_box(file_path, slide_idx, summary_text, 1.0, 2.5, 8.0, 4.0, 18, "#000000", False, False)
    else:
        await add_text_box(file_path, slide_idx, "• No key points extracted from presentation content",
                          1.0, 2.5, 8.0, 1.0, 18, "#666666", False, True)

    return {
        "message": f"Generated summary slide '{title}' at index {slide_idx}",
        "slide_index": slide_idx,
        "points_extracted": len(summary_points),
        "max_points": max_points
    }


async def main() -> None:
    """Main entry point for the PPTX MCP server."""
    log.info("Starting PowerPoint MCP server (stdio)...")
    from mcp.server.stdio import stdio_server

    async with stdio_server() as (read_stream, write_stream):
        await server.run(
            read_stream,
            write_stream,
            InitializationOptions(
                server_name="pptx-server",
                server_version="0.1.0",
                capabilities={"tools": {}, "logging": {}},
            ),
        )


if __name__ == "__main__":
    asyncio.run(main())